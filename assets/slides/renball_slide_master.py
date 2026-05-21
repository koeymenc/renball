"""
renball_slide_master.py — Shared PowerPoint slide builder for Renball
=====================================================================

The canonical Renball slide module. Every per-project `code/build_slides.py`
imports from here and never reimplements styling.

Design (mirrors CLAUDE.md §5):
- 16:9 aspect, 13.333" × 7.5"
- Strict slide zones:
    Title area:  0.4"  – 1.3"  from top   (eyebrow + serif title + accent line)
    Body / charts: 1.6" – 6.8"  from top   (5.2" of body height)
    Footer:      6.9"  – 7.4"  from top   (logo bottom-left, page number bottom-right)
- Background:        #080c08
- Headlines:         Instrument Serif, #e8efe8 (accent runs in italic #c1ff72)
- Body:              DM Sans, #8fa68f
- Captions / labels: JetBrains Mono uppercase, #5a735a
- Chart palette:     matplotlib equivalents (MPL_*) mirror the page tokens

Title accent runs
-----------------
Wrap a phrase in asterisks to render it in italic + accent. Example:

    add_title_slide(prs, "Not All Goals Are *Equal*.")

splits the title into two non-accent runs and one italic-accent run.
The parser is exposed as `parse_accent_runs(text)` for callers that want
to do their own rendering.

API surface
-----------
Zone constants:
    TITLE_AREA_TOP/BOTTOM, BODY_TOP, BODY_BOTTOM, BODY_LEFT, BODY_MAX_W,
    CAPTION_TOP, FOOTER_TOP, SLIDE_W, SLIDE_H

Tokens (re-exported):
    BG, BG_ELEVATED, BG_CARD, SURFACE, BORDER, BORDER_LIGHT,
    TEXT, TEXT_SECONDARY, TEXT_MUTED, ACCENT, ACCENT_DIM,
    SERIF, SANS, MONO,
    MPL_BG, MPL_BG_CARD, MPL_BORDER, MPL_BORDER_LIGHT,
    MPL_TEXT, MPL_TEXT_SEC, MPL_TEXT_MUTED, MPL_ACCENT, MPL_SECONDARY

Deck:
    create_deck() -> Presentation
    finalize_deck(prs)              # stamps page numbers, skips title + closing

High-level slide builders:
    add_title_slide(prs, title, *, subtitle=None, date=None, project_num=None)
    add_section_slide(prs, section_num, section_title)
    add_bullet_slide(prs, title, bullets, *, caption=None, eyebrow=None)
    add_two_column_slide(prs, title, left, right, *, eyebrow=None, caption=None)
    add_chart_slide(prs, title, fig, *, caption=None, eyebrow=None)
                                    # auto-fits chart inside body zone
    add_table_slide(prs, title, headers, rows, *, caption=None, eyebrow=None,
                    highlight_first_row=False, col_widths=None)
    add_quote_slide(prs, quote, *, attribution=None)
    add_closing_slide(prs, *, headline="Thanks.", url="renball.com",
                      project_path=None, lines=None)

Low-level primitives (for bespoke layouts):
    add_blank_slide(prs) -> slide                # bg + logo, no page number
    add_logo(slide)
    add_page_number(slide, n, total=None)
    add_section_header(slide, eyebrow, title)
    textbox / add_run / add_text / add_mono_label / add_rect / add_hline
    new_chart(width=11.0, height=4.2)
    chart_to_bytes(fig)
    add_chart_picture(slide, fig, x, y, w, h=None)   # h optional; if omitted,
                                                     # height auto-scales from
                                                     # figure aspect

All x/y/w/h arguments are in inches.
"""

from __future__ import annotations

import io
import re
from typing import Optional, Sequence

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ════════════════════════════════════════════════════════════════════════
# Design tokens (mirror assets/css/design-tokens.css)
# ════════════════════════════════════════════════════════════════════════
BG           = RGBColor(0x08, 0x0C, 0x08)
BG_ELEVATED  = RGBColor(0x0E, 0x15, 0x0E)
BG_CARD      = RGBColor(0x11, 0x1A, 0x11)
SURFACE      = RGBColor(0x1A, 0x26, 0x1A)
BORDER       = RGBColor(0x1E, 0x2E, 0x1E)
BORDER_LIGHT = RGBColor(0x2D, 0x3D, 0x2D)
TEXT           = RGBColor(0xE8, 0xEF, 0xE8)
TEXT_SECONDARY = RGBColor(0x8F, 0xA6, 0x8F)
TEXT_MUTED     = RGBColor(0x5A, 0x73, 0x5A)
ACCENT     = RGBColor(0xC1, 0xFF, 0x72)
ACCENT_DIM = RGBColor(0x8F, 0xBF, 0x45)

SERIF = "Instrument Serif"
SANS  = "DM Sans"
MONO  = "JetBrains Mono"

# Matplotlib equivalents (#rrggbb)
MPL_BG           = "#080c08"
MPL_BG_CARD      = "#111a11"
MPL_BORDER       = "#1e2e1e"
MPL_BORDER_LIGHT = "#2d3d2d"
MPL_TEXT         = "#e8efe8"
MPL_TEXT_SEC     = "#8fa68f"
MPL_TEXT_MUTED   = "#5a735a"
MPL_ACCENT       = "#c1ff72"
MPL_SECONDARY    = "#72c1ff"

# ════════════════════════════════════════════════════════════════════════
# Slide geometry (16:9) — strict zones
# ════════════════════════════════════════════════════════════════════════
SLIDE_W = 13.333  # inches
SLIDE_H = 7.5     # inches

# Title zone (eyebrow + serif title + accent line)
TITLE_AREA_TOP    = 0.4
TITLE_AREA_BOTTOM = 1.3

# Body / chart zone — everything between header and footer
BODY_TOP     = 1.6
BODY_BOTTOM  = 6.8
BODY_LEFT    = 0.7
BODY_MAX_W   = 12.0

# Caption sits inside the body zone so it never collides with the footer
CAPTION_TOP    = 6.5
CAPTION_BOTTOM = 6.8
CAPTION_GAP    = 0.15   # gap between chart bottom and caption top

# Footer (logo bottom-left + page number bottom-right)
FOOTER_TOP    = 6.9
FOOTER_BOTTOM = 7.4
FOOTER_Y      = 6.95    # baseline y for both logo + page number

# Internal flag attribute name used by finalize_deck() to skip slides that
# shouldn't get a page number (title, closing).
_NO_PAGENUM_ATTR = "_renball_no_pagenum"


# ════════════════════════════════════════════════════════════════════════
# Accent-run parser  (*...* markers → italic + accent runs)
# ════════════════════════════════════════════════════════════════════════
_ACCENT_RUN_RE = re.compile(r"\*([^*]+)\*")


def parse_accent_runs(text: str) -> list[tuple[str, bool]]:
    """
    Split a title string on *...* markers.

    Returns a list of (segment, is_accent) tuples. Asterisks are stripped
    from the segments. Adjacent text runs are emitted as (..., False).

        >>> parse_accent_runs("Not All Goals Are *Equal*.")
        [('Not All Goals Are ', False), ('Equal', True), ('.', False)]
    """
    parts: list[tuple[str, bool]] = []
    pos = 0
    for m in _ACCENT_RUN_RE.finditer(text):
        if m.start() > pos:
            parts.append((text[pos:m.start()], False))
        parts.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        parts.append((text[pos:], False))
    if not parts:
        parts = [(text, False)]
    return parts


def _render_title_runs(paragraph, title: str, *,
                       font: str = SERIF, size: int = 32,
                       base_color: RGBColor = TEXT,
                       accent_color: RGBColor = ACCENT,
                       accent_italic: bool = True) -> None:
    """Render `title` (with optional *...* accent markers) into one paragraph."""
    for segment, is_accent in parse_accent_runs(title):
        if is_accent:
            add_run(paragraph, segment, font=font, size=size,
                    color=accent_color, italic=accent_italic)
        else:
            add_run(paragraph, segment, font=font, size=size,
                    color=base_color)


# ════════════════════════════════════════════════════════════════════════
# Deck creation + finalisation
# ════════════════════════════════════════════════════════════════════════
def create_deck() -> Presentation:
    """Create a new 16:9 Renball-sized Presentation."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def finalize_deck(prs: Presentation) -> None:
    """
    Stamp page numbers on every slide except those tagged with
    _NO_PAGENUM_ATTR (title and closing slides). Call once after all
    slides have been added.
    """
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if getattr(slide, _NO_PAGENUM_ATTR, False):
            continue
        add_page_number(slide, i, total)


# ════════════════════════════════════════════════════════════════════════
# Low-level primitives
# ════════════════════════════════════════════════════════════════════════
def add_blank_slide(prs: Presentation):
    """Add a blank slide with Renball background + bottom-left wordmark.
    No page number — finalize_deck() adds those later."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    add_logo(slide)
    return slide


def add_logo(slide):
    """Bottom-left ren·ball wordmark ('ball' in accent), 16pt."""
    tb = slide.shapes.add_textbox(
        Inches(0.45), Inches(FOOTER_Y),
        Inches(2), Inches(0.4),
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = "ren"
    r1.font.name = SERIF; r1.font.size = Pt(16); r1.font.color.rgb = TEXT
    r2 = p.add_run(); r2.text = "ball"
    r2.font.name = SERIF; r2.font.size = Pt(16); r2.font.color.rgb = ACCENT
    return tb


def add_page_number(slide, n: int, total: Optional[int] = None) -> None:
    """Bottom-right page number in muted mono. '07' or '07 / 13' format."""
    text = f"{n:02d} / {total:02d}" if total else f"{n:02d}"
    tb = slide.shapes.add_textbox(
        Inches(SLIDE_W - 1.4), Inches(FOOTER_Y),
        Inches(1.2), Inches(0.35),
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.name = MONO; r.font.size = Pt(9)
    r.font.color.rgb = TEXT_MUTED


def textbox(slide, x, y, w, h, *,
            anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    """Construct a zero-margin textbox at (x, y) of size (w, h) inches."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.word_wrap = wrap
    tf.paragraphs[0].alignment = align
    return tb


def add_run(paragraph, text, *, font=SANS, size=14, color=TEXT,
            bold=False, italic=False):
    """Append a Run with EXPLICIT font.name (no PPTX-default reliance)."""
    r = paragraph.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return r


def add_text(slide, x, y, w, h, text, *, font=SANS, size=14, color=TEXT,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = textbox(slide, x, y, w, h, anchor=anchor, align=align)
    add_run(tb.text_frame.paragraphs[0], text,
            font=font, size=size, color=color, bold=bold, italic=italic)
    return tb


def add_mono_label(slide, x, y, w, h, text, *,
                   color=TEXT_MUTED, size=9, align=PP_ALIGN.LEFT):
    """Uppercase mono label (eyebrows, captions, table headers)."""
    return add_text(slide, x, y, w, h, text.upper(),
                    font=MONO, size=size, color=color, align=align)


def add_rect(slide, x, y, w, h, fill, *, border=None, border_pt=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_pt or 0.75)
    else:
        shp.line.fill.background()
    shp.text_frame.text = ""
    shp.text_frame.margin_left = shp.text_frame.margin_right = 0
    shp.text_frame.margin_top = shp.text_frame.margin_bottom = 0
    return shp


def add_hline(slide, x, y, w, *, color=BORDER, pt=1):
    """Horizontal line of pt thickness."""
    return add_rect(slide, x, y, w, pt / 72, color)


def add_section_header(slide, eyebrow: Optional[str], title: str, *,
                       accent_line: bool = True,
                       title_size: int = 32) -> None:
    """
    Standard Renball slide-top pattern, strictly inside the title area
    (0.4"–1.3" from top):
        eyebrow (mono 10pt)      at y=0.50, ends ~0.72
        serif title (32pt)        at y=0.78, ends ~1.32
        accent underline          at y=1.40   (just below title area)

    Title supports *...* markers for italic + accent runs.
    """
    if eyebrow:
        add_mono_label(slide, BODY_LEFT, 0.50, BODY_MAX_W, 0.25, eyebrow,
                       color=ACCENT, size=10)
    tb = textbox(slide, BODY_LEFT, 0.78, BODY_MAX_W, 0.55, align=PP_ALIGN.LEFT)
    _render_title_runs(tb.text_frame.paragraphs[0], title,
                       font=SERIF, size=title_size,
                       base_color=TEXT, accent_color=ACCENT)
    if accent_line:
        add_hline(slide, BODY_LEFT, 1.40, 1.2, color=ACCENT, pt=2)


# ════════════════════════════════════════════════════════════════════════
# Matplotlib helpers
# ════════════════════════════════════════════════════════════════════════
def new_chart(width: float = 11.0, height: float = 4.2):
    """
    Pre-styled (fig, ax) — Renball palette, no top/right spines.
    Defaults sized so 12" wide → ≤ 4.6" tall in PPTX, fitting comfortably
    inside the body zone (BODY_TOP→CAPTION_TOP) with room for a caption.
    """
    fig, ax = plt.subplots(figsize=(width, height), facecolor=MPL_BG)
    ax.set_facecolor(MPL_BG)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    for side in ("bottom", "left"):
        ax.spines[side].set_edgecolor(MPL_BORDER)
    ax.tick_params(colors=MPL_TEXT_MUTED, labelsize=10)
    ax.grid(axis="y", color=MPL_BORDER, alpha=0.7, linewidth=0.6)
    ax.set_axisbelow(True)
    return fig, ax


def chart_to_bytes(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200,
                facecolor=MPL_BG, bbox_inches="tight", pad_inches=0.2)
    buf.seek(0)
    plt.close(fig)
    return buf


def add_chart_picture(slide, fig, x: float, y: float,
                      w: float, h: Optional[float] = None):
    """
    Render `fig` to PNG and place it at (x, y).

    If `h` is given, the image is forced to (w × h) — aspect ratio of the
    underlying figure is not preserved. If `h` is omitted, the image
    auto-scales by the figure's aspect ratio (the original behaviour).
    """
    buf = chart_to_bytes(fig)
    if h is None:
        return slide.shapes.add_picture(buf, Inches(x), Inches(y),
                                        width=Inches(w))
    return slide.shapes.add_picture(buf, Inches(x), Inches(y),
                                    width=Inches(w), height=Inches(h))


def _fit_in_body(fig, *, max_w: float, max_h: float):
    """
    Compute the PPTX placement for a matplotlib figure inside a bounding
    box of (max_w × max_h) inches. Returns (chart_w, chart_h) preserving
    figure aspect ratio.
    """
    w_inch, h_inch = fig.get_size_inches()
    ratio = w_inch / max(h_inch, 1e-6)
    if max_w / ratio <= max_h:
        chart_w = max_w
        chart_h = chart_w / ratio
    else:
        chart_h = max_h
        chart_w = chart_h * ratio
    return chart_w, chart_h


# ════════════════════════════════════════════════════════════════════════
# High-level slide builders
# ════════════════════════════════════════════════════════════════════════
def add_title_slide(prs: Presentation, title: str, *,
                    subtitle: Optional[str] = None,
                    date: Optional[str] = None,
                    project_num: Optional[str] = None):
    """
    Deck opener. Big serif title supports *...* accent markers (italic +
    accent colour on the marked phrase). Title slide skips page numbering.
    """
    s = add_blank_slide(prs)
    setattr(s, _NO_PAGENUM_ATTR, True)

    if project_num:
        add_mono_label(s, BODY_LEFT, 2.0, BODY_MAX_W, 0.3,
                       f"RESEARCH · {project_num}",
                       color=ACCENT, size=11)

    tb = textbox(s, BODY_LEFT, 2.5, BODY_MAX_W, 2.6,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True)
    _render_title_runs(tb.text_frame.paragraphs[0], title,
                       font=SERIF, size=54,
                       base_color=TEXT, accent_color=ACCENT)

    if subtitle:
        add_mono_label(s, BODY_LEFT, 5.2, BODY_MAX_W, 0.4, subtitle,
                       color=TEXT_SECONDARY, size=12)
    if date:
        add_mono_label(s, BODY_LEFT, 5.7, BODY_MAX_W, 0.4, date,
                       color=TEXT_MUTED, size=10)
    add_hline(s, BODY_LEFT, 6.15, 1.2, color=ACCENT, pt=2)
    return s


def add_section_slide(prs: Presentation,
                      section_num: str, section_title: str):
    """Section divider: large accent number + serif title (supports *...*)."""
    s = add_blank_slide(prs)
    add_text(s, BODY_LEFT, 2.6, 4.0, 2.0, section_num,
             font=SERIF, size=120, color=ACCENT)
    tb = textbox(s, 4.7, 3.2, 8.0, 2.2,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    _render_title_runs(tb.text_frame.paragraphs[0], section_title,
                       font=SERIF, size=44,
                       base_color=TEXT, accent_color=ACCENT)
    add_hline(s, 4.7, 5.3, 1.5, color=ACCENT, pt=2)
    return s


def add_bullet_slide(prs: Presentation, title: str,
                     bullets: Sequence[str], *,
                     caption: Optional[str] = None,
                     eyebrow: Optional[str] = None):
    """Section header + 3–6 bullets (em-dash markers, accent)."""
    s = add_blank_slide(prs)
    add_section_header(s, eyebrow, title)

    n = max(len(bullets), 1)
    top    = BODY_TOP
    bottom = CAPTION_TOP - CAPTION_GAP if caption else BODY_BOTTOM
    row_h  = min(0.85, max(0.45, (bottom - top) / n))
    y = top
    for bullet in bullets:
        add_text(s, BODY_LEFT, y, 0.4, row_h, "—",
                 font=SERIF, size=22, color=ACCENT)
        add_text(s, BODY_LEFT + 0.45, y + 0.05, BODY_MAX_W - 0.5, row_h, bullet,
                 font=SANS, size=15, color=TEXT_SECONDARY)
        y += row_h + 0.15

    if caption:
        add_text(s, BODY_LEFT, CAPTION_TOP, BODY_MAX_W, 0.3, caption,
                 font=MONO, size=10, color=TEXT_MUTED)
    return s


def add_two_column_slide(prs: Presentation, title: str,
                         left: dict, right: dict, *,
                         eyebrow: Optional[str] = None,
                         caption: Optional[str] = None):
    """
    Section header + two side-by-side cards.
    Each column dict supports keys: 'label' (mono eyebrow), 'title' (serif),
    'body' (sans), 'accent' (bool — highlight the card border).
    """
    s = add_blank_slide(prs)
    add_section_header(s, eyebrow, title)

    card_w = (BODY_MAX_W - 0.25) / 2
    card_h = (CAPTION_TOP - CAPTION_GAP if caption else BODY_BOTTOM) - BODY_TOP
    cards = [
        (BODY_LEFT, left),
        (BODY_LEFT + card_w + 0.25, right),
    ]
    for x, content in cards:
        accent = bool(content.get("accent"))
        add_rect(s, x, BODY_TOP, card_w, card_h, BG_CARD,
                 border=ACCENT if accent else BORDER_LIGHT,
                 border_pt=1.5 if accent else 0.75)
        inner_x = x + 0.25
        inner_w = card_w - 0.5
        if content.get("label"):
            add_mono_label(s, inner_x, BODY_TOP + 0.2, inner_w, 0.3,
                           content["label"],
                           color=ACCENT if accent else TEXT_SECONDARY,
                           size=10)
        if content.get("title"):
            add_text(s, inner_x, BODY_TOP + 0.55, inner_w, 0.8,
                     content["title"], font=SERIF, size=24, color=TEXT)
        if content.get("body"):
            add_text(s, inner_x, BODY_TOP + 1.50, inner_w, card_h - 1.7,
                     content["body"], font=SANS, size=14,
                     color=TEXT_SECONDARY)
    if caption:
        add_text(s, BODY_LEFT, CAPTION_TOP, BODY_MAX_W, 0.3, caption,
                 font=MONO, size=10, color=TEXT_MUTED)
    return s


def add_chart_slide(prs: Presentation, title: str, fig, *,
                    caption: Optional[str] = None,
                    eyebrow: Optional[str] = None):
    """
    Section header + auto-fitted matplotlib figure + optional mono caption.

    The chart is placed inside the body zone (BODY_TOP–BODY_BOTTOM, or
    BODY_TOP–(CAPTION_TOP − CAPTION_GAP) when a caption is present). The
    figure's aspect ratio is preserved — height is clamped to the body
    zone so the chart never overflows the slide.
    """
    s = add_blank_slide(prs)
    add_section_header(s, eyebrow, title)

    chart_bottom = (CAPTION_TOP - CAPTION_GAP) if caption else BODY_BOTTOM
    max_w = BODY_MAX_W
    max_h = chart_bottom - BODY_TOP

    chart_w, chart_h = _fit_in_body(fig, max_w=max_w, max_h=max_h)

    # Centre the chart horizontally + vertically inside the body zone
    chart_x = BODY_LEFT + (max_w - chart_w) / 2
    chart_y = BODY_TOP + (max_h - chart_h) / 2

    add_chart_picture(s, fig, chart_x, chart_y, chart_w, chart_h)

    if caption:
        add_text(s, BODY_LEFT, CAPTION_TOP, BODY_MAX_W, 0.3, caption,
                 font=MONO, size=10, color=TEXT_MUTED)
    return s


def add_table_slide(prs: Presentation, title: str,
                    headers: Sequence[str],
                    rows: Sequence[Sequence], *,
                    caption: Optional[str] = None,
                    eyebrow: Optional[str] = None,
                    highlight_first_row: bool = False,
                    col_widths: Optional[Sequence[float]] = None):
    """
    Section header + plain table (rendered as text rows so colours stay
    on-brand). Sum of col_widths is clamped to BODY_MAX_W.
    """
    s = add_blank_slide(prs)
    add_section_header(s, eyebrow, title)

    n_cols = len(headers)
    if col_widths is None:
        col_widths = [BODY_MAX_W / n_cols] * n_cols
    elif sum(col_widths) > BODY_MAX_W:
        scale = BODY_MAX_W / sum(col_widths)
        col_widths = [w * scale for w in col_widths]

    x_positions = [BODY_LEFT]
    for w in col_widths[:-1]:
        x_positions.append(x_positions[-1] + w)

    # Header row
    y = BODY_TOP
    add_hline(s, BODY_LEFT, y - 0.05, sum(col_widths),
              color=BORDER_LIGHT, pt=1)
    for i, hd in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        add_mono_label(s, x_positions[i], y, col_widths[i] - 0.1, 0.4,
                       str(hd), color=TEXT_MUTED, size=10, align=align)
    y += 0.5
    add_hline(s, BODY_LEFT, y - 0.05, sum(col_widths), color=BORDER, pt=1)

    # Body rows — split available height across rows
    bottom = (CAPTION_TOP - CAPTION_GAP) if caption else BODY_BOTTOM
    available_h = bottom - y
    n_rows = max(len(rows), 1)
    row_h = min(0.5, max(0.32, available_h / n_rows - 0.05))

    for ri, row in enumerate(rows):
        is_highlight = highlight_first_row and ri == 0
        for ci, cell in enumerate(row):
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            font = SANS if ci == 0 else MONO
            color = TEXT if ci == 0 else TEXT_SECONDARY
            size = 13 if ci == 0 else 12
            if is_highlight and ci > 0:
                color = ACCENT
            bold = is_highlight
            add_text(s, x_positions[ci], y, col_widths[ci] - 0.1, row_h,
                     str(cell), font=font, size=size, color=color,
                     align=align, bold=bold)
        y += row_h + 0.05

    if caption:
        add_text(s, BODY_LEFT, CAPTION_TOP, BODY_MAX_W, 0.3, caption,
                 font=MONO, size=10, color=TEXT_MUTED)
    return s


def add_quote_slide(prs: Presentation, quote: str, *,
                    attribution: Optional[str] = None):
    """Single large centred pull-quote, italic, accent leading mark."""
    s = add_blank_slide(prs)
    add_text(s, 1.5, 0.7, 10.5, 0.6, "“",
             font=SERIF, size=60, color=ACCENT, align=PP_ALIGN.LEFT)
    add_text(s, 1.5, 2.4, 10.5, 3.5, quote,
             font=SERIF, size=34, color=TEXT, italic=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if attribution:
        add_mono_label(s, 1.5, 6.0, 10.5, 0.4, f"— {attribution}",
                       color=TEXT_MUTED, size=11)
    add_hline(s, 1.5, 6.55, 1.5, color=ACCENT, pt=2)
    return s


def add_closing_slide(prs: Presentation, *,
                      headline: str = "Thanks.",
                      url: str = "renball.com",
                      project_path: Optional[str] = None,
                      lines: Optional[Sequence[str]] = None):
    """
    Final slide: large serif sign-off + URL in accent. No page number.
    `project_path` is a mono subtitle (URL slug). `lines` is an optional
    list of mono reference lines.
    """
    s = add_blank_slide(prs)
    setattr(s, _NO_PAGENUM_ATTR, True)

    add_text(s, BODY_LEFT, 1.4, BODY_MAX_W, 1.0, headline,
             font=SERIF, size=42, color=TEXT)
    add_hline(s, BODY_LEFT, 2.5, 1.2, color=ACCENT, pt=2)

    add_text(s, BODY_LEFT, 3.1, BODY_MAX_W, 0.5,
             "Full methodology, source code and data:",
             font=SANS, size=15, color=TEXT_SECONDARY)
    add_text(s, BODY_LEFT, 3.7, BODY_MAX_W, 1.0, url,
             font=SERIF, size=56, color=ACCENT)

    if project_path:
        add_mono_label(s, BODY_LEFT, 5.0, BODY_MAX_W, 0.4, project_path,
                       color=TEXT_SECONDARY, size=12)

    if lines:
        # Dynamic spacing — fit n lines between y=5.4 and y=FOOTER_TOP − 0.2
        # so the last line stops cleanly above the logo at FOOTER_Y.
        avail_top = 5.4
        avail_bottom = FOOTER_TOP - 0.2  # 6.70
        n = max(len(lines), 1)
        step = (avail_bottom - avail_top) / n
        y = avail_top
        for line in lines:
            add_text(s, BODY_LEFT, y, BODY_MAX_W, step, line,
                     font=MONO, size=12, color=TEXT_MUTED)
            y += step
    return s
