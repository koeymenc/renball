"""
build_slides.py — Renball Project 01 slide builder
===================================================

Generates a 12-slide 16:9 deck for the penalty conversion project. All
numbers are pulled from ../data/data.json so the deck stays in sync
with the build.

Styling and slide layouts come exclusively from
`assets/slides/renball_slide_master.py` — never reimplement them here.

Run:
    python code/build_slides.py
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

# ── Pull in the shared Renball slide master ─────────────────
SCRIPT_DIR  = Path(__file__).resolve().parent
PROJECT_DIR = SCRIPT_DIR.parent
REPO_ROOT   = PROJECT_DIR.parent.parent
SHARED      = REPO_ROOT / "assets" / "slides"
sys.path.insert(0, str(SHARED))

import renball_slide_master as R  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Pt           # noqa: E402

DATA_JSON   = PROJECT_DIR / "data" / "data.json"
OUTPUT_PPTX = PROJECT_DIR / "slides.pptx"


# ════════════════════════════════════════════════════════════════════════
# Slide content
# ════════════════════════════════════════════════════════════════════════
def slide_title(prs, d):
    date_str = (d.get("generated_at") or "")[:7] or "May 2026"
    R.add_title_slide(
        prs,
        title="Should the Fouled Player Take the Penalty *Himself*?",
        subtitle="A ONE-TAILED Z-TEST · TOP 5 EUROPEAN LEAGUES",
        date=date_str,
        project_num="01",
    )


def slide_tldr(prs, d):
    t = d["primary"]["treatment"]
    c = d["primary"]["control"]
    tst = d["primary"]["test"]

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "TL;DR", "Three *sentences*")

    rows = [
        ("THE QUESTION",
         "When a goalkeeper fouls a player who then takes the penalty himself, "
         "does conversion drop?"),
        ("THE TEST",
         f"One-tailed two-proportion z-test. Treatment n = {t['attempts']} "
         f"({t['rate']}%). Control n = {c['attempts']:,} ({c['rate']}%)."),
        ("THE RESULT",
         f"Δ = {tst['rate_diff_pp']:+.2f} pp (predicted direction) but "
         f"p = {tst['p_value']:.3f}, one-tailed — fail to reject H₀ at the available power."),
    ]
    y = R.BODY_TOP
    for label, body in rows:
        R.add_mono_label(s, R.BODY_LEFT, y, 2.4, 0.3, label,
                         color=R.ACCENT, size=10)
        R.add_text(s, R.BODY_LEFT + 2.5, y - 0.07, 9.3, 1.2, body,
                   font=R.SERIF, size=20, color=R.TEXT)
        y += 1.55


def slide_theory(prs, d):
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "THE THEORY", "Why we'd *expect* a drop")
    R.add_text(s, R.BODY_LEFT, R.BODY_TOP, R.BODY_MAX_W, 0.7,
               "The just-fouled player walks straight to the spot.",
               font=R.SERIF, size=24, color=R.TEXT, italic=True)
    R.add_text(s, R.BODY_LEFT, R.BODY_TOP + 1.0, 6.0, 3.4,
               "Theory of the taker: physically jarred — body contact, hit the ground, "
               "adrenaline. Emotionally agitated, narrowed focus.",
               font=R.SANS, size=15, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT + 6.3, R.BODY_TOP + 1.0, 5.7, 3.4,
               "Theory of the keeper: locked in. The foul that gave the penalty was his fault. "
               "Nothing left to lose. Maximum motivation.",
               font=R.SANS, size=15, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "→ If both pull conversion down, the bench should hand the ball to a designated specialist.",
               font=R.SANS, size=13, color=R.ACCENT, italic=True)


def slide_data(prs, d):
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "DATA", "What we built the test on")

    seasons = d["metadata"]["seasons"]
    seasons_label = f"{seasons[0]} → {seasons[-1]}" if seasons else "—"
    t_n = d["primary"]["treatment"]["attempts"]

    stats = [
        (f"{d['overall']['attempts']:,}",            "PENALTIES (LOOSE FILTER)"),
        (f"{d['metadata']['n_matches_strict']:,}",   "SINGLE-PENALTY MATCHES"),
        (f"{t_n:,}",                                  "TREATMENT-CELL N"),
        (str(len(d["metadata"]["leagues"])),         "LEAGUES"),
    ]
    x = R.BODY_LEFT
    for value, label in stats:
        R.add_rect(s, x, R.BODY_TOP, 2.85, 1.6, R.BG_CARD,
                   border=R.BORDER, border_pt=0.75)
        R.add_text(s, x + 0.2, R.BODY_TOP + 0.15, 2.6, 0.9, value,
                   font=R.SERIF, size=42, color=R.ACCENT)
        R.add_mono_label(s, x + 0.2, R.BODY_TOP + 1.05, 2.6, 0.3, label,
                         color=R.TEXT_MUTED, size=9)
        x += 3.0

    y = R.BODY_TOP + 1.95   # 3.55
    R.add_text(s, R.BODY_LEFT, y, R.BODY_MAX_W, 0.4,
               f"Seasons in scope: {seasons_label}",
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, y + 0.5, R.BODY_MAX_W, 0.4,
               "Competitions: Premier League, La Liga, Serie A, Bundesliga, Ligue 1.",
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, y + 1.0, R.BODY_MAX_W, 0.4,
               "Source: FBref player-match data, consolidated locally as FULL_DICT_ROWS.pkl.",
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Strict filter = matches with at most one penalty event "
               "(unambiguous GK-caused / taker-was-fouled attribution).",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_method(prs, d):
    R.add_two_column_slide(
        prs,
        title="One-tailed two-proportion z-test",
        left={
            "label": "TREATMENT",
            "title": "GK fouled the taker",
            "body": "Matches where the goalkeeper conceded the foul AND the fouled "
                    "player is the same player who took the penalty.",
            "accent": True,
        },
        right={
            "label": "CONTROL",
            "title": "All other penalties",
            "body": "Every other strict-filter match — GK fouled but a different player took, "
                    "or outfielder caused the foul (any taker).",
        },
        eyebrow="METHOD",
        caption="H₀: conversion_treatment = conversion_control   ·   "
                "H₁: conversion_treatment < conversion_control (one-tailed, α = 0.05)   ·   "
                "95% Wilson CIs · statsmodels.proportions_ztest",
    )


def slide_by_league(prs, d):
    leagues = d["by_league"]
    labels = [l["label"] for l in leagues]
    rates = [l["rate"] for l in leagues]
    n_vals = [l["attempts"] for l in leagues]

    fig, ax = R.new_chart(width=11, height=4.3)
    bars = ax.bar(labels, rates, color=R.MPL_ACCENT,
                  width=0.55, edgecolor="none")
    ax.set_ylim(0, 100)
    ax.set_ylabel("Conversion rate (%)",
                  color=R.MPL_TEXT_SEC, fontsize=11)
    for bar, val, n in zip(bars, rates, n_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, val + 1.2,
                f"{val:.1f}%",
                ha="center", color=R.MPL_TEXT, fontsize=10)
        ax.text(bar.get_x() + bar.get_width() / 2, -6,
                f"n = {n:,}",
                ha="center", color=R.MPL_TEXT_MUTED, fontsize=8)

    R.add_chart_slide(prs,
                      "Conversion rate by league",
                      fig,
                      caption="Across the top 5 leagues conversion clusters tightly in the high 70s — "
                              "league of origin is not a major source of variation.",
                      eyebrow="CONTEXT")


def slide_by_season(prs, d):
    seasons = d["by_season"]
    labels = [s["season"] for s in seasons]
    rates = [s["rate"] for s in seasons]
    n_vals = [s["attempts"] for s in seasons]

    fig, ax = R.new_chart(width=11, height=4.3)
    ax.bar(labels, rates, color=R.MPL_ACCENT,
           width=0.55, edgecolor="none")
    ax.set_ylim(0, 100)
    ax.set_ylabel("Conversion rate (%)",
                  color=R.MPL_TEXT_SEC, fontsize=11)
    for i, (val, n) in enumerate(zip(rates, n_vals)):
        ax.text(i, val + 1.2, f"{val:.1f}%",
                ha="center", color=R.MPL_TEXT, fontsize=9)
        ax.text(i, -6, f"n = {n:,}",
                ha="center", color=R.MPL_TEXT_MUTED, fontsize=8)
    import matplotlib.pyplot as plt
    plt.setp(ax.get_xticklabels(), rotation=30, ha="right")

    R.add_chart_slide(prs,
                      "Conversion across seasons",
                      fig,
                      caption="Year-on-year stability is high — the ~77–80% band holds across the period in scope.",
                      eyebrow="CONTEXT")


def slide_primary(prs, d):
    """Bespoke layout: split chart on the left + result card on the right."""
    t = d["primary"]["treatment"]
    c = d["primary"]["control"]
    tst = d["primary"]["test"]

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "PRIMARY TEST · ONE-TAILED",
                         "Treatment vs *control*")

    body_top = R.BODY_TOP
    card_h = R.CAPTION_TOP - R.CAPTION_GAP - body_top   # 4.75"

    # Left chart — explicit (w, h) so we control placement precisely
    chart_w = 7.4
    chart_h = card_h
    fig, ax = R.new_chart(width=chart_w, height=chart_h)
    labels = ["Treatment\n(GK fouled the taker)", "Control\n(all other penalties)"]
    rates = [t["rate"], c["rate"]]
    bars = ax.bar(labels, rates, color=[R.MPL_ACCENT, R.MPL_SECONDARY],
                  width=0.5, edgecolor="none")
    cis = [(t["ci_lo"], t["ci_hi"]), (c["ci_lo"], c["ci_hi"])]
    for i, (lo, hi) in enumerate(cis):
        x = bars[i].get_x() + bars[i].get_width() / 2
        ax.plot([x, x], [lo, hi], color=R.MPL_TEXT, linewidth=1.4)
        ax.plot([x - 0.07, x + 0.07], [lo, lo], color=R.MPL_TEXT, linewidth=1.4)
        ax.plot([x - 0.07, x + 0.07], [hi, hi], color=R.MPL_TEXT, linewidth=1.4)
    for bar, val in zip(bars, rates):
        ax.text(bar.get_x() + bar.get_width() / 2, val - 6,
                f"{val:.2f}%",
                ha="center", color=R.MPL_BG, fontsize=14, fontweight="bold")
    ax.set_ylim(50, 95)
    ax.set_ylabel("Conversion rate (%)",
                  color=R.MPL_TEXT_SEC, fontsize=11)
    R.add_chart_picture(s, fig, R.BODY_LEFT, body_top, chart_w, chart_h)

    # Right card
    x0 = R.BODY_LEFT + chart_w + 0.3   # 8.4
    card_w = R.BODY_MAX_W - chart_w - 0.3   # 4.3"
    R.add_rect(s, x0, body_top, card_w, card_h, R.BG_CARD,
               border=R.ACCENT, border_pt=1.5)

    R.add_mono_label(s, x0 + 0.25, body_top + 0.2, card_w - 0.5, 0.3,
                     "TREATMENT", color=R.TEXT_MUTED, size=9)
    R.add_text(s, x0 + 0.25, body_top + 0.5, card_w - 0.5, 0.7,
               f"{t['rate']:.2f}%",
               font=R.SERIF, size=28, color=R.ACCENT)
    R.add_text(s, x0 + 0.25, body_top + 1.15, card_w - 0.5, 0.3,
               f"n = {t['attempts']:,}  ·  95% CI [{t['ci_lo']:.2f}%, {t['ci_hi']:.2f}%]",
               font=R.MONO, size=9, color=R.TEXT_SECONDARY)

    R.add_mono_label(s, x0 + 0.25, body_top + 1.65, card_w - 0.5, 0.3,
                     "CONTROL", color=R.TEXT_MUTED, size=9)
    R.add_text(s, x0 + 0.25, body_top + 1.95, card_w - 0.5, 0.7,
               f"{c['rate']:.2f}%",
               font=R.SERIF, size=28, color=R.ACCENT)
    R.add_text(s, x0 + 0.25, body_top + 2.6, card_w - 0.5, 0.3,
               f"n = {c['attempts']:,}  ·  95% CI [{c['ci_lo']:.2f}%, {c['ci_hi']:.2f}%]",
               font=R.MONO, size=9, color=R.TEXT_SECONDARY)

    R.add_hline(s, x0 + 0.25, body_top + 3.15, card_w - 0.5,
                color=R.BORDER_LIGHT, pt=1)

    R.add_text(s, x0 + 0.25, body_top + 3.30, card_w - 0.5, 0.45,
               f"Δ = {tst['rate_diff_pp']:+.2f} pp   z = {tst['z_stat']:+.3f}",
               font=R.MONO, size=11, color=R.TEXT)
    verdict = "Reject H₀" if tst["significant_at_05"] else "Fail to reject H₀"
    R.add_text(s, x0 + 0.25, body_top + 3.75, card_w - 0.5, 0.4,
               f"p = {tst['p_value']:.4f} (one-tailed)  ·  {verdict}",
               font=R.MONO, size=11, color=R.ACCENT)

    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Whiskers: 95% Wilson confidence intervals. "
               "Treatment CI is wide — n is the binding constraint.",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_matrix(prs, d):
    """Bespoke 2×2 matrix layout."""
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "2×2 CONTEXT", "Where the *treatment cell* sits")

    m = d["matrix"]
    by_rc = {(cc["row"], cc["col"]): cc for cc in m["cells"]}

    # Grid: 1 header row + 2 data rows. Total grid height fits within body zone.
    # header 0.7 + 2 × 1.85 = 4.4". Body has 5.2" available before caption.
    label_w = 2.5
    cell_w = 4.0
    cell_h = 1.85
    grid_w = label_w + 2 * cell_w   # 10.5"
    grid_x = R.BODY_LEFT + (R.BODY_MAX_W - grid_w) / 2  # centre horizontally
    grid_y = R.BODY_TOP

    # Column headers
    R.add_rect(s, grid_x, grid_y, label_w, 0.7, R.BG_ELEVATED, border=R.BORDER)
    for ci, col in enumerate(m["cols"]):
        cx = grid_x + label_w + ci * cell_w
        R.add_rect(s, cx, grid_y, cell_w, 0.7, R.BG_ELEVATED, border=R.BORDER)
        R.add_mono_label(s, cx + 0.1, grid_y + 0.2, cell_w - 0.2, 0.3,
                         col["label"], color=R.TEXT_SECONDARY, size=10,
                         align=PP_ALIGN.CENTER)

    # Body
    for ri, row in enumerate(m["rows"]):
        ry = grid_y + 0.7 + ri * cell_h
        R.add_rect(s, grid_x, ry, label_w, cell_h, R.BG_ELEVATED, border=R.BORDER)
        R.add_mono_label(s, grid_x + 0.15, ry + cell_h / 2 - 0.2, label_w - 0.3, 0.5,
                         row["label"], color=R.TEXT_SECONDARY, size=10)

        for ci, col in enumerate(m["cols"]):
            cx = grid_x + label_w + ci * cell_w
            cell = by_rc[(row["key"], col["key"])]
            shp = R.add_rect(s, cx, ry, cell_w, cell_h, R.BG_CARD, border=R.BORDER)
            if cell.get("is_treatment"):
                shp.line.color.rgb = R.ACCENT
                shp.line.width = Pt(2)
                R.add_rect(s, cx + 0.05, ry + 0.05, cell_w - 0.1, cell_h - 0.1, R.SURFACE)
                R.add_mono_label(s, cx + 0.15, ry + 0.15, cell_w - 0.3, 0.3,
                                 "TREATMENT", color=R.ACCENT, size=8,
                                 align=PP_ALIGN.CENTER)
            R.add_text(s, cx + 0.1, ry + 0.45, cell_w - 0.2, 0.7,
                       f"{cell['rate']:.2f}%",
                       font=R.SERIF, size=30, color=R.ACCENT,
                       align=PP_ALIGN.CENTER)
            R.add_mono_label(s, cx + 0.1, ry + 1.15, cell_w - 0.2, 0.3,
                             f"n = {cell['attempts']:,}",
                             color=R.TEXT_MUTED, size=10,
                             align=PP_ALIGN.CENTER)

    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Accent-bordered cell = treatment. The other three pool into control.",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_per_league(prs, d):
    """Bespoke wide table with sig coloring + low-n italic."""
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "PER-LEAGUE REPLICATION",
                         "The same test, *league by league*")

    blocks = d["primary"]["by_league"]

    headers = ["LEAGUE", "T. RATE", "T. N", "C. RATE", "C. N", "Z", "P (1-T)", "SIG"]
    col_w = [2.6, 1.4, 1.0, 1.4, 1.2, 1.2, 1.4, 1.4]
    col_x = [R.BODY_LEFT]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    header_y = R.BODY_TOP
    R.add_hline(s, R.BODY_LEFT, header_y - 0.05, sum(col_w),
                color=R.BORDER_LIGHT, pt=1)
    for i, hd in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        R.add_mono_label(s, col_x[i], header_y, col_w[i] - 0.1, 0.4,
                         hd, color=R.TEXT_MUTED, size=10, align=align)
    R.add_hline(s, R.BODY_LEFT, header_y + 0.45, sum(col_w),
                color=R.BORDER, pt=1)

    y = header_y + 0.6
    for blk in blocks:
        tt = blk["treatment"]; cc = blk["control"]; tst = blk["test"]
        low_n = tt["attempts"] < 30
        z_str = f"{tst['z_stat']:+.3f}" if tst["z_stat"] is not None else "—"
        p_str = f"{tst['p_value']:.4f}" if tst["p_value"] is not None else "—"
        sig_str = "Yes" if tst["significant_at_05"] else "No"
        sig_color = R.ACCENT if tst["significant_at_05"] else R.TEXT_MUTED
        t_rate_color = R.TEXT_MUTED if low_n else R.TEXT

        R.add_text(s, col_x[0], y, col_w[0] - 0.1, 0.4, blk["label"],
                   font=R.SANS, size=13, color=R.TEXT)
        R.add_text(s, col_x[1], y, col_w[1] - 0.1, 0.4, f"{tt['rate']:.2f}%",
                   font=R.MONO, size=12, color=t_rate_color,
                   align=PP_ALIGN.RIGHT, italic=low_n)
        R.add_text(s, col_x[2], y, col_w[2] - 0.1, 0.4, f"{tt['attempts']}",
                   font=R.MONO, size=12, color=t_rate_color,
                   align=PP_ALIGN.RIGHT, italic=low_n)
        R.add_text(s, col_x[3], y, col_w[3] - 0.1, 0.4, f"{cc['rate']:.2f}%",
                   font=R.MONO, size=12, color=R.TEXT_SECONDARY,
                   align=PP_ALIGN.RIGHT)
        R.add_text(s, col_x[4], y, col_w[4] - 0.1, 0.4, f"{cc['attempts']:,}",
                   font=R.MONO, size=12, color=R.TEXT_SECONDARY,
                   align=PP_ALIGN.RIGHT)
        R.add_text(s, col_x[5], y, col_w[5] - 0.1, 0.4, z_str,
                   font=R.MONO, size=12, color=R.TEXT,
                   align=PP_ALIGN.RIGHT)
        R.add_text(s, col_x[6], y, col_w[6] - 0.1, 0.4, p_str,
                   font=R.MONO, size=12, color=R.TEXT,
                   align=PP_ALIGN.RIGHT)
        sig_tag = sig_str + (" · low n" if low_n else "")
        R.add_text(s, col_x[7], y, col_w[7] - 0.1, 0.4, sig_tag,
                   font=R.MONO, size=12, color=sig_color,
                   align=PP_ALIGN.RIGHT)
        y += 0.5

    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Per-league treatment cells of 13–25 leave each test severely underpowered — "
               "pooled is the only inference to take seriously.",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_conclusion(prs, d):
    tst = d["primary"]["test"]
    t = d["primary"]["treatment"]
    c = d["primary"]["control"]
    diff_abs = abs(tst["rate_diff_pp"])
    direction = "lower" if tst["rate_diff_pp"] < 0 else "higher"

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "CONCLUSION", "What the data *says*")

    # Big headline with accent emphasis
    tb = R.textbox(s, R.BODY_LEFT, R.BODY_TOP, R.BODY_MAX_W, 1.0)
    p = tb.text_frame.paragraphs[0]
    R.add_run(p, "The bench answer: ",
              font=R.SERIF, size=32, color=R.TEXT)
    R.add_run(p, "can't yet tell.",
              font=R.SERIF, size=32, color=R.ACCENT, italic=True)

    R.add_text(s, R.BODY_LEFT, R.BODY_TOP + 1.2, R.BODY_MAX_W, 2.0,
               f"Conversion is {diff_abs:.2f} pp {direction} when the fouled player takes the kick himself "
               f"({t['rate']:.2f}% vs {c['rate']:.2f}% for everyone else). The point estimate moves in the "
               "predicted direction, but the treatment cell is too small to rule chance out "
               f"(p = {tst['p_value']:.3f}, one-tailed).",
               font=R.SANS, size=15, color=R.TEXT_SECONDARY)

    R.add_text(s, R.BODY_LEFT, R.BODY_TOP + 3.4, R.BODY_MAX_W, 1.2,
               "→ No statistical warrant to override the fouled player who wants to take it himself. "
               "Equally, no clean evidence he's at full strength.",
               font=R.SANS, size=13, color=R.ACCENT, italic=True)


def slide_sources(prs, d):
    R.add_closing_slide(
        prs,
        headline="Full methodology, source code and data:",
        url="renball.com",
        project_path="/projects/01-penalty-analysis/",
        lines=[
            "→ methodology      formal H₀ / H₁, test assumptions, limitations",
            "→ analysis.html    interactive dashboard, live numbers from data.json",
            "→ code/build.py    Python pipeline — raw pickles → aggregated JSON",
        ],
    )


# ════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════
def main() -> None:
    with open(DATA_JSON, encoding="utf-8") as f:
        d = json.load(f)

    prs = R.create_deck()
    slide_title(prs, d)        # 1
    slide_tldr(prs, d)         # 2
    slide_theory(prs, d)       # 3
    slide_data(prs, d)         # 4
    slide_method(prs, d)       # 5
    slide_by_league(prs, d)    # 6
    slide_by_season(prs, d)    # 7
    slide_primary(prs, d)      # 8
    slide_matrix(prs, d)       # 9
    slide_per_league(prs, d)   # 10
    slide_conclusion(prs, d)   # 11
    slide_sources(prs, d)      # 12
    R.finalize_deck(prs)

    prs.save(OUTPUT_PPTX)
    size_kb = OUTPUT_PPTX.stat().st_size / 1024
    print(f"[OK] Wrote {OUTPUT_PPTX.relative_to(PROJECT_DIR)} "
          f"({size_kb:.1f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
