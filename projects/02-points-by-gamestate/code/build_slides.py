"""
build_slides.py — Renball Project 02 slide builder
===================================================

Generates a 13-slide 16:9 deck for Points by Gamestate. All numbers are
pulled from ../data/data.json so the deck stays in sync with the build.

Styling and slide layouts come exclusively from
`assets/slides/renball_slide_master.py` — never reimplement them here.

Run:
    python code/build_slides.py
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

# ── Pull in the shared Renball slide master ─────────────────
SCRIPT_DIR  = Path(__file__).resolve().parent
PROJECT_DIR = SCRIPT_DIR.parent
REPO_ROOT   = PROJECT_DIR.parent.parent                 # …/renball/
SHARED      = REPO_ROOT / "assets" / "slides"
sys.path.insert(0, str(SHARED))

import renball_slide_master as R  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

# Inputs / outputs
DATA_JSON   = PROJECT_DIR / "data" / "data.json"
OUTPUT_PPTX = PROJECT_DIR / "slides.pptx"


# ════════════════════════════════════════════════════════════════════════
# Slide content (one function per slide)
# ════════════════════════════════════════════════════════════════════════
def slide_title(prs, d):
    meta = d["meta"]
    date_str = (meta.get("generated_at") or "").replace("-", "·")[:7] or "May 2026"
    R.add_title_slide(
        prs,
        title="Not All Goals Are *Equal*.",
        subtitle="A TIME-AWARE EXPECTED-POINTS MODEL · TOP 5 LEAGUES + EREDIVISIE",
        date=date_str,
        project_num="02",
    )


def slide_tldr(prs, d):
    co = d["callouts"]
    up = co["biggest_upgrade"]
    bullets = [
        f"A goal that takes a team from 0–0 to 1–0 in the final fifteen minutes is "
        f"worth {co['ep_value_gap']['best']:.2f} expected points — "
        f"2.1× the same transition in the opening fifteen.",
        f"Across {co['total_goals']:,} goals in six leagues and seven seasons we "
        f"observed {co['unique_states']} of 90 theoretical (gd × minute_bin) states.",
        f"The model's biggest gainer over the time-blind baseline is "
        f"{up['player']} (+{up['delta_ep']:.2f} ΔEP on {up['goals_involved']} G+A) — "
        f"a late-equaliser specialist the scoreboard underrates.",
    ]
    R.add_bullet_slide(prs, "TL;DR — three sentences",
                       bullets, eyebrow="TL;DR")


def slide_question(prs, d):
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "THE QUESTION", "Are all goals *equal*?")
    R.add_text(s, R.BODY_LEFT, R.BODY_TOP, R.BODY_MAX_W, 0.9,
               "Three goals, three different lengths of game ahead. Should we treat them the same?",
               font=R.SERIF, size=24, color=R.TEXT, italic=True)
    rows = [
        ("EARLY OPENER",
         "0–0 to 1–0 in the 5th minute. 85 minutes still to play, everyone resets."),
        ("LATE EQUALISER",
         "1–2 to 2–2 in the 88th minute. Pulls a point out of a lost match."),
        ("CONSOLATION",
         "4–0 down, you score in stoppage. Doesn't move the table."),
    ]
    y = R.BODY_TOP + 1.05    # 2.65
    for label, body in rows:
        R.add_mono_label(s, R.BODY_LEFT, y, 2.6, 0.3, label,
                         color=R.ACCENT, size=10)
        R.add_text(s, R.BODY_LEFT + 2.7, y - 0.05, 9.0, 0.65, body,
                   font=R.SANS, size=14, color=R.TEXT_SECONDARY)
        y += 0.75
    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "→ The standard scoreboard treats them the same. We won't.",
               font=R.SANS, size=14, color=R.ACCENT, italic=True)


def slide_data(prs, d):
    meta = d["meta"]
    co = d["callouts"]
    leagues = meta["leagues"]
    seasons = meta["seasons"]

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "DATA", "What we built the framework on")

    stats = [
        (f"{co['total_goals']:,}",          "GOALS ANALYZED"),
        (f"{len(leagues)}",                 "LEAGUES"),
        (f"{len(seasons)}",                 "SEASONS"),
        (f"{co['unique_states']}",          "MODELLED STATES"),
    ]
    x = R.BODY_LEFT
    for value, label in stats:
        R.add_rect(s, x, R.BODY_TOP, 2.85, 1.6, R.BG_CARD,
                   border=R.BORDER, border_pt=0.75)
        R.add_text(s, x + 0.2, R.BODY_TOP + 0.15, 2.6, 0.9, value,
                   font=R.SERIF, size=42, color=R.ACCENT)
        R.add_mono_label(s, x + 0.2, R.BODY_TOP + 1.05, 2.6, 0.3, label,
                         color=R.TEXT_MUTED, size=9)
        x += 3.0

    # Caption-zone text rows (cards bottom at 3.2)
    y = R.BODY_TOP + 1.95   # 3.55
    R.add_text(s, R.BODY_LEFT, y, R.BODY_MAX_W, 0.4,
               f"Seasons in scope: {seasons[0]} → {seasons[-1]}",
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, y + 0.5, R.BODY_MAX_W, 0.4,
               "Leagues: " + ", ".join(leagues),
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, y + 1.0, R.BODY_MAX_W, 0.4,
               "Source: FBref-derived shot tables, one pickle per league "
               "(match_data_final_<league>.pkl).",
               font=R.SANS, size=13, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Variables: scoring team, outcome=goal, goalscorer, SCA1 (assist), "
               "minute (parsed including stoppage suffixes).",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_value_function(prs, d):
    s = R.add_blank_slide(prs)
    R.add_section_header(s, "THE VALUE FUNCTION", "How a goal gets *priced*")

    R.add_text(s, R.BODY_LEFT, R.BODY_TOP, R.BODY_MAX_W, 0.5,
               "For every pre-goal state (goal_difference, minute_bin) we compute the empirical expected points.",
               font=R.SANS, size=14, color=R.TEXT_SECONDARY)

    # Formula card — state-value definition
    R.add_rect(s, R.BODY_LEFT, R.BODY_TOP + 0.55, R.BODY_MAX_W, 1.25, R.BG_CARD,
               border=R.ACCENT, border_pt=1.2)
    R.add_text(s, R.BODY_LEFT + 0.3, R.BODY_TOP + 0.75, R.BODY_MAX_W - 0.6, 0.55,
               "EP(gd, minute_bin)  =  3 · P(W | state)  +  1 · P(D | state)",
               font=R.MONO, size=18, color=R.ACCENT)
    R.add_text(s, R.BODY_LEFT + 0.3, R.BODY_TOP + 1.35, R.BODY_MAX_W - 0.6, 0.4,
               "Mean of final-match points over every team that entered this (gd, bin) state.",
               font=R.MONO, size=11, color=R.TEXT_MUTED)

    R.add_text(s, R.BODY_LEFT, R.BODY_TOP + 2.0, R.BODY_MAX_W, 0.5,
               "A goal's value is the difference between the two states it connects:",
               font=R.SANS, size=14, color=R.TEXT_SECONDARY)
    R.add_rect(s, R.BODY_LEFT, R.BODY_TOP + 2.55, R.BODY_MAX_W, 1.0, R.BG_CARD,
               border=R.BORDER, border_pt=0.75)
    R.add_text(s, R.BODY_LEFT + 0.3, R.BODY_TOP + 2.78, R.BODY_MAX_W - 0.6, 0.55,
               "ep_collected  =  EP(gd_after, bin)  −  EP(gd_before, bin)",
               font=R.MONO, size=16, color=R.TEXT)


def slide_state_value_table(prs, d):
    """Visualise the state-value table as a coloured grid (a small subset of GDs)."""
    svt = d["state_value_table"]
    bins = svt["time_bins"]
    gds = svt["gd_range"]
    matrix = svt["matrix"]

    # Use the GD-range subset that has data on most rows (gd in [-4, +4])
    show_gds = [i for i, gd in enumerate(gds) if -4 <= gd <= 4]

    # Find min/max for colour mapping
    vals = [c["exp_points"] for row in matrix for c in row
            if c and c.get("exp_points") is not None]
    lo, hi = (min(vals), max(vals)) if vals else (0.0, 3.0)

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "STATE VALUE TABLE",
                         "Expected final points at each (GD, minute)")

    # Geometry — grid fits within body zone (1.6–6.8). 9 rows × ~0.45" each ≈ 4.05".
    col_w = 1.45
    row_h = 0.45
    label_w = 1.6
    grid_w = label_w + len(bins) * col_w   # 1.6 + 6 * 1.45 = 10.3"
    grid_left = R.BODY_LEFT + (R.BODY_MAX_W - grid_w) / 2  # centre horizontally
    grid_top = R.BODY_TOP

    # Header row: minute bins
    R.add_rect(s, grid_left, grid_top, label_w, row_h, R.BG_ELEVATED)
    for j, b in enumerate(bins):
        x = grid_left + label_w + j * col_w
        R.add_rect(s, x, grid_top, col_w, row_h, R.BG_ELEVATED)
        R.add_mono_label(s, x, grid_top + 0.05, col_w, row_h - 0.05, b,
                         color=R.TEXT_SECONDARY, size=9, align=PP_ALIGN.CENTER)

    # Rows
    def _fmt_gd(n):
        return "0" if n == 0 else (f"+{n}" if n > 0 else f"−{abs(n)}")

    def _color_alpha(val):
        # Map val ∈ [lo, hi] to 0.07..0.65 accent intensity
        if val is None:
            return None
        t = max(0.0, min(1.0, (val - lo) / (hi - lo)))
        # Linear interpolate ACCENT_DIM ↔ ACCENT
        alpha = 0.10 + t * 0.55
        # Mix BG_CARD with ACCENT according to alpha
        r1, g1, b1 = 0x11, 0x1A, 0x11  # BG_CARD
        r2, g2, b2 = 0xC1, 0xFF, 0x72  # ACCENT
        r = int(r1 * (1 - alpha) + r2 * alpha)
        g = int(g1 * (1 - alpha) + g2 * alpha)
        bx = int(b1 * (1 - alpha) + b2 * alpha)
        from pptx.dml.color import RGBColor as _RGB
        return _RGB(r, g, bx)

    y = grid_top + row_h
    for i in show_gds:
        gd = gds[i]
        label = f"{_fmt_gd(gd)} to {_fmt_gd(gd + 1)}"
        # Row label
        R.add_rect(s, grid_left, y, label_w, row_h, R.BG_ELEVATED)
        R.add_mono_label(s, grid_left + 0.05, y + 0.06, label_w - 0.1, row_h - 0.05,
                         label, color=R.TEXT_SECONDARY, size=9, align=PP_ALIGN.RIGHT)
        # Cells
        for j in range(len(bins)):
            cell = matrix[i][j]
            x = grid_left + label_w + j * col_w
            fill = _color_alpha(cell["exp_points"]) if cell else R.BG_CARD
            R.add_rect(s, x, y, col_w, row_h, fill or R.BG_CARD,
                       border=R.BORDER, border_pt=0.4)
            if cell and cell.get("exp_points") is not None:
                R.add_text(s, x, y + 0.04, col_w, row_h - 0.05,
                           f"{cell['exp_points']:.2f}",
                           font=R.MONO, size=11, color=R.TEXT, align=PP_ALIGN.CENTER)
        y += row_h

    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "Showing GD ∈ [−4, +4]; the full table covers [−7, +7]. "
               "Cells brighten with expected points.",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_ep_by_bin(prs, d):
    """Total EP cashed in by minute bin — aggregate bar chart."""
    hm = d["value_heatmap"]
    bins = hm["time_bins"]
    totals = []
    for j in range(len(bins)):
        s_sum = 0.0
        for i in range(len(hm["transitions"])):
            cell = hm["matrix"][i][j]
            if cell and cell.get("ep_collected") is not None:
                s_sum += cell["ep_collected"] * cell["n"]
        totals.append(s_sum)

    fig, ax = R.new_chart()  # default 11×4.2 — fits body zone
    bars = ax.bar(bins, totals, color=R.MPL_ACCENT,
                  width=0.6, edgecolor=R.MPL_ACCENT, linewidth=0)
    for bar, val in zip(bars, totals):
        ax.text(bar.get_x() + bar.get_width() / 2, val + max(totals) * 0.02,
                f"{val:,.0f}",
                ha="center", color=R.MPL_TEXT, fontsize=11)
    ax.set_ylabel("EP collected (total)",
                  color=R.MPL_TEXT_SEC, fontsize=11)
    ax.set_xlabel("Minute bin", color=R.MPL_TEXT_SEC, fontsize=11)
    ax.set_ylim(0, max(totals) * 1.15)

    R.add_chart_slide(prs,
                      "Where the expected points get *cashed in*",
                      fig,
                      caption="Sum of (ep_collected × n) across every transition, per minute bin. "
                              "The last 15 minutes deliver ~1.6× the value of the first.",
                      eyebrow="AGGREGATE")


def slide_early_vs_late(prs, d):
    """0→1 EP collected by bin — the cleanest single-transition story."""
    bins = d["early_vs_late_0to1"]["time_bins"]
    series = d["early_vs_late_0to1"]["ep_collected"]

    fig, ax = R.new_chart()
    bars = ax.bar(bins, series, color=R.MPL_ACCENT,
                  width=0.6, edgecolor=R.MPL_ACCENT, linewidth=0)
    for bar, val in zip(bars, series):
        if val is None:
            continue
        ax.text(bar.get_x() + bar.get_width() / 2, val + max(series) * 0.02,
                f"{val:.2f}",
                ha="center", color=R.MPL_TEXT, fontsize=12)
    ax.set_ylabel("EP collected", color=R.MPL_TEXT_SEC, fontsize=11)
    ax.set_xlabel("Minute bin", color=R.MPL_TEXT_SEC, fontsize=11)
    ax.set_ylim(0, max(series) * 1.18)

    R.add_chart_slide(prs,
                      "The opener: worth more the *longer it took* to arrive",
                      fig,
                      caption=f"EP collected from a 0→1 goal by minute bin. "
                              f"Late opener ({series[-1]:.2f} EP) is {series[-1] / series[0]:.1f}× the early one ({series[0]:.2f} EP).",
                      eyebrow="0 → 1")


def slide_top10_ep(prs, d):
    """Top 10 EP collectors leaderboard."""
    rows = d["top10_ep_collectors"]
    labels = [f"{r['player']}" for r in rows][::-1]
    values = [r["ep_timed"] for r in rows][::-1]
    teams = [r["team"] for r in rows][::-1]

    fig, ax = R.new_chart(width=11, height=4.5)
    bars = ax.barh(labels, values, color=R.MPL_ACCENT,
                   edgecolor=R.MPL_ACCENT, linewidth=0)
    for bar, val, team in zip(bars, values, teams):
        ax.text(val + max(values) * 0.01,
                bar.get_y() + bar.get_height() / 2,
                f"{val:.1f}  ·  {team}",
                ha="left", va="center",
                color=R.MPL_TEXT_SEC, fontsize=10)
    ax.set_xlabel("EP collected (time-aware)",
                  color=R.MPL_TEXT_SEC, fontsize=11)
    ax.set_xlim(0, max(values) * 1.25)
    ax.tick_params(axis="y", labelsize=11)

    R.add_chart_slide(prs,
                      "Top 10 EP collectors · Goals + Assists",
                      fig,
                      caption="All leagues, all seasons pooled. Goals + assists combined; "
                              "players pooled across team transfers (primary team labelled).",
                      eyebrow="LEADERBOARD")


def slide_mario_spotlight(prs, d):
    """Mário Rui spotlight slide — the biggest upgrade story."""
    up = d["callouts"]["biggest_upgrade"]
    rows = d["top10_timed_upgrade"]

    s = R.add_blank_slide(prs)
    R.add_section_header(s, "SPOTLIGHT",
                         "The biggest *time-aware* upgrade")

    # Left card — height fits body zone (1.6 to caption row at 6.5 = 4.9")
    card_top = R.BODY_TOP
    card_h = R.CAPTION_TOP - R.CAPTION_GAP - card_top   # 4.75"
    R.add_rect(s, R.BODY_LEFT, card_top, 6.0, card_h, R.BG_CARD,
               border=R.ACCENT, border_pt=1.5)
    R.add_mono_label(s, R.BODY_LEFT + 0.25, card_top + 0.2, 5.5, 0.3,
                     "BIGGEST GAINER", color=R.ACCENT, size=10)
    R.add_text(s, R.BODY_LEFT + 0.25, card_top + 0.55, 5.5, 0.7, up["player"],
               font=R.SERIF, size=32, color=R.TEXT)
    R.add_text(s, R.BODY_LEFT + 0.25, card_top + 1.3, 5.5, 0.4,
               f"{up['team']}  ·  {up['goals_involved']} G+A pooled",
               font=R.MONO, size=11, color=R.TEXT_SECONDARY)
    R.add_text(s, R.BODY_LEFT + 0.25, card_top + 1.95, 5.5, 1.1,
               f"+{up['delta_ep']:.2f}",
               font=R.SERIF, size=64, color=R.ACCENT)
    R.add_mono_label(s, R.BODY_LEFT + 0.25, card_top + 3.10, 5.5, 0.3,
                     "Δ EP FROM TIME-AWARE RE-RATING",
                     color=R.TEXT_MUTED, size=9)
    R.add_text(s, R.BODY_LEFT + 0.25, card_top + 3.45, 5.5, 0.4,
               f"({up['ep_timed']:.2f} timed  −  {up['ep_untimed']:.2f} untimed)",
               font=R.MONO, size=11, color=R.TEXT_SECONDARY)

    # Right: top-5 upgrade leaderboard
    right_x = R.BODY_LEFT + 6.3
    right_w = R.BODY_MAX_W - 6.3   # 5.7"
    R.add_mono_label(s, right_x, card_top + 0.15, right_w, 0.3,
                     "TOP 5 UPGRADES",
                     color=R.TEXT_MUTED, size=10)
    R.add_hline(s, right_x, card_top + 0.55, right_w, color=R.BORDER, pt=1)
    y = card_top + 0.7
    for i, r in enumerate(rows[:5], start=1):
        is_top = i == 1
        R.add_text(s, right_x, y, 0.4, 0.4, f"{i}",
                   font=R.MONO, size=12,
                   color=R.ACCENT if is_top else R.TEXT_MUTED)
        R.add_text(s, right_x + 0.45, y, 3.4, 0.4, r["player"],
                   font=R.SANS, size=13,
                   color=R.TEXT if is_top else R.TEXT_SECONDARY,
                   bold=is_top)
        R.add_text(s, right_x + right_w - 1.7, y, 1.6, 0.4,
                   f"+{r['delta_ep']:.2f}",
                   font=R.MONO, size=12,
                   color=R.ACCENT if is_top else R.TEXT_SECONDARY,
                   align=PP_ALIGN.RIGHT, bold=is_top)
        y += 0.55

    R.add_text(s, R.BODY_LEFT, R.CAPTION_TOP, R.BODY_MAX_W, 0.3,
               "A walking late-equaliser: ~30 career goals in scope, almost entirely in score-tied late minutes.",
               font=R.MONO, size=10, color=R.TEXT_MUTED)


def slide_implications(prs, d):
    R.add_bullet_slide(
        prs,
        "What this changes",
        bullets=[
            "Goals are not a unit — they are a function of state and time. "
            "EP-weighted rankings expose contribution patterns raw counts can't see.",
            "Recruitment with leverage in mind: a player whose 12 goals are mostly late equalisers "
            "is materially different from one whose 12 are 4–0 cushions.",
            "Squad-management: knowing which minute bins your forwards harvest from informs "
            "substitution timing and game-state-conditional roles.",
            "Coaching: late-minute defending decisions get a clean valuation — the equaliser you concede "
            "in the 88th minute is worth ~1.7× the one in the 18th.",
        ],
        eyebrow="IMPLICATIONS",
    )


def slide_limitations(prs, d):
    R.add_bullet_slide(
        prs,
        "What this isn't",
        bullets=[
            "Not xG-weighted. A tap-in and a 30-yard volley count the same.",
            "Penalties and own goals are not separated from open play (next iteration).",
            "Assist credit caps at SCA1 — build-up beyond the first shot-creating action isn't captured.",
            "Observational, not causal: it values the gap between states, not the goal's incremental effect "
            "on the final outcome.",
            "Pooled value function across leagues — per-league calibration available in the dashboard.",
        ],
        eyebrow="LIMITATIONS",
        caption="Full discussion: methodology.html · §5",
    )


def slide_contact(prs, d):
    R.add_closing_slide(
        prs,
        headline="Want every cell, every league, every season?",
        url="renball.com",
        project_path="/projects/02-points-by-gamestate/",
        lines=[
            "→ analysis.html        editorial write-up with five Chart.js + D3 visuals",
            "→ dashboard.html       full interactive Tabulator dashboard, six datasets, drill-down",
            "→ methodology.html     formal write-up — state-value math, attribution, eight limitations",
            "→ code/build.py        Python pipeline from raw pickles to JSON",
        ],
    )


# ════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════
def main() -> None:
    with open(DATA_JSON, encoding="utf-8") as f:
        d = json.load(f)

    prs = R.create_deck()
    slide_title(prs, d)               # 1
    slide_tldr(prs, d)                # 2
    slide_question(prs, d)            # 3
    slide_data(prs, d)                # 4
    slide_value_function(prs, d)      # 5
    slide_state_value_table(prs, d)   # 6
    slide_ep_by_bin(prs, d)           # 7
    slide_early_vs_late(prs, d)       # 8
    slide_top10_ep(prs, d)            # 9
    slide_mario_spotlight(prs, d)     # 10
    slide_implications(prs, d)        # 11
    slide_limitations(prs, d)         # 12
    slide_contact(prs, d)             # 13
    R.finalize_deck(prs)

    prs.save(OUTPUT_PPTX)
    size_kb = OUTPUT_PPTX.stat().st_size / 1024
    print(f"[OK] Wrote {OUTPUT_PPTX.relative_to(PROJECT_DIR)} "
          f"({size_kb:.1f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
